import sklearn.utils
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.model_selection import train_test_split
from sklearn.naive_bayes import MultinomialNB
from sklearn.metrics import accuracy_score
from packaging import version
from tkinter import Tk, filedialog, StringVar, Label, Entry, Button, Frame, messagebox, Toplevel, OptionMenu
from pptx import Presentation
from docx import Document
import jieba, json, os, joblib, PyPDF2, winreg, sys, requests, threading


# 读取pptx文件
def pptx(name):
    prs = Presentation(str(name))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs


# 读取docx文件
def docx(name):
    document = Document(str(name))
    ps = [paragraph.text for paragraph in document.paragraphs]
    return ps


# 读取pdf文件
def pdf(name):
    mypdf = open(name, mode='rb')
    pdf_document = PyPDF2.PdfReader(mypdf)
    an = []
    for i in range(len(pdf_document.pages)):
        t = pdf_document.pages[i].extract_text().split(' ')
        if len(t) > 1:
            an += t
    return an


# 对读取的数据进行预处理
def structure(lst, md=2):
    def chdl(r):
        r = jieba.lcut(r)
        i = 0
        while i < len(r):
            if len(r[i]) < 2:
                r.pop(i)
                continue
            i += 1
        return r

    ls = []
    for i in lst:
        if len(i) < 2:
            continue
        cnt = [0] * 2
        x = [''] * 2
        i = i.lower()
        for s in i:
            if '\u4e00' <= s <= '\u9fff':
                cnt[0] += 1
                x[0] += s
            elif 'a' <= s <= 'z':
                cnt[1] += 1
                x[1] += s
            else:
                if len(x[0]) > 1: ls += chdl(x[0])
                if len(x[1]) > 2: ls.append(x[1])
                x = [''] * 2
        if len(x[0]) > 1: ls += chdl(x[0])
        if len(x[1]) > 2: ls.append(x[1])
    if md == 2:
        return ' '.join(ls)
    dic = {}
    for i in ls:
        if i in dic:
            dic[i] += 1
        else:
            dic[i] = 1
    return dic


# 通过文件名调用对应函数
def getcont(fn, md=2):
    typ = fn.split('.')[-1]
    if typ == 'pptx':
        a = pptx(fn)
    elif typ == 'docx':
        a = docx(fn)
    elif typ == 'pdf':
        a = pdf(fn)
    else:
        return False
    return structure(a, md)


# 模型的训练与保存
def model_sv(rec):
    doc = []
    labels = []
    
    # 创建进度对话框
    progress_window = Toplevel()
    progress_window.title("模型训练进度")
    window_width = 400
    window_height = 100
    screen_width = progress_window.winfo_screenwidth()
    screen_height = progress_window.winfo_screenheight()
    x = (screen_width/2) - (window_width/2)
    y = (screen_height/2) - (window_height/2)
    progress_window.geometry('%dx%d+%d+%d' % (window_width, window_height, x, y))
    
    # 设置进度窗口始终在主窗口之上
    progress_window.transient()  # 设置为临时窗口
    progress_window.grab_set()   # 捕获所有事件
    
    # 进度条标签
    outer_label = Label(progress_window, text="准备预处理数据...")
    outer_label.pack(pady=5)
    
    inner_label = Label(progress_window, text="")
    inner_label.pack(pady=5)
    
    # 进度条
    from tkinter.ttk import Progressbar
    progress_bar = Progressbar(progress_window, orient="horizontal", length=300, mode="determinate")
    progress_bar.pack(pady=10)

    print('开始预处理数据')
    total_outer = len(rec)
    count_outer = 0
    
    # 更新总进度条最大值（分两个阶段：预处理和训练）
    total_steps = total_outer + 1  # 类别处理 + 模型训练
    current_step = 0
    
    for i in rec:
        count_outer += 1
        outer_label.config(text=f"处理类别 {i} [{count_outer}/{total_outer}]")
        progress_window.update()
        
        total_inner = len(rec[i])
        count_inner = 0
        
        # 设置当前进度条为文件处理进度
        progress_bar["maximum"] = total_inner
        progress_bar["value"] = 0
        
        for j in rec[i]:
            count_inner += 1
            try:
                inner_label.config(text=f"处理文件 [{count_inner}/{total_inner}]")
                progress_bar["value"] = count_inner
                progress_window.update()
                
                doc.append(getcont(j))
                labels.append(i)
            except Exception as e:
                print(f'\n添加数据失败：{j}，错误信息：{str(e)}')
        
        current_step += 1
        progress_bar["maximum"] = total_steps
        progress_bar["value"] = current_step
    
    outer_label.config(text="开始训练模型")
    inner_label.config(text="")
    progress_bar["value"] = 0
    progress_window.update()
    
    print('\n开始训练模型')
    vectorizer = TfidfVectorizer()
    X = vectorizer.fit_transform(doc)
    X_train, X_test, y_train, y_test = train_test_split(X, labels, test_size=0.3, random_state=42)
    
    clf = MultinomialNB()
    clf.fit(X_train, y_train)
    
    y_pred = clf.predict(X_test)
    accuracy = accuracy_score(y_test, y_pred)
    print(f"模型准确率: {accuracy}")
    
    # 更新进度条表示训练完成
    progress_bar["maximum"] = 100
    progress_bar["value"] = 100
    outer_label.config(text="模型训练完成")
    inner_label.config(text=f"准确率: {accuracy:.2%}")
    progress_window.update()
    
    # 使用新的图形界面获取模型名称
    model_name = ask_model_name(tmp['modph'])
    
    if model_name:
        t = os.path.join(tmp['modph'], model_name)
        vt, cf = os.path.join(t, 'vectorizer.pkl'), os.path.join(t, 'classifier.pkl')
        
        if not os.path.exists(t):
            os.makedirs(t)
            
        try:
            joblib.dump(vectorizer, vt)
            joblib.dump(clf, cf)
            print('模型保存成功')
            
            # 显示完成窗口一段时间
            progress_window.after(1000, progress_window.destroy)
            return True
        except Exception as e:
            print(f'保存模型时发生错误：{str(e)}')
    
    print('模型训练和保存已取消')
    progress_window.destroy()
    return False

# 新增的图形界面函数
def ask_model_name(model_path):
    """通过GUI对话框获取模型名称"""
    import tkinter as tk
    from tkinter import simpledialog
    
    class ModelNameDialog(simpledialog.Dialog):
        def __init__(self, parent, title=None):
            self.existing_models = [d for d in os.listdir(model_path) if os.path.isdir(os.path.join(model_path, d))]
            super().__init__(parent, title)
            
        def body(self, master):
            Label(master, text="请输入模型名称:").grid(row=0, sticky="w")
            self.name_entry = Entry(master)
            self.name_entry.grid(row=1, sticky="ew")
            
            if self.existing_models:
                Label(master, text="已存在的模型:").grid(row=2, sticky="w", pady=(10, 0))
                self.model_listbox = tk.Listbox(master, height=5)
                self.model_listbox.grid(row=3, sticky="ew")
                
                for model in self.existing_models:
                    self.model_listbox.insert(tk.END, model)
                    
            return self.name_entry
            
        def validate(self):
            name = self.name_entry.get().strip()
            if not name:
                messagebox.showerror("错误", "模型名称不能为空！")
                return 0
                
            if not name.isalnum() and not all(c in name for c in "-_"):
                messagebox.showerror("错误", "模型名称只能包含字母、数字、下划线和连字符！")
                return 0
                
            if name in self.existing_models:
                messagebox.showerror("错误", f"模型 '{name}' 已存在，请选择其他名称！")
                return 0
                
            self.result = name
            return 1
            
    # 创建临时Tk实例进行对话
    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口
    
    try:
        dialog = ModelNameDialog(root, "命名模型")
        return dialog.result
    finally:
        root.destroy()


# 模型的使用
def model_use():
    f = []
    lis = []
    for i in os.listdir(config['collect']):
        if not os.path.isdir(i) and i.split('.')[-1] in ['pptx', 'docx', 'pdf']:
            try:
                t = getcont(os.path.join(config['collect'], i))
                if t != False and len(t.split()) > 15:  # 跳过文本内容过少的文档，至少10个词
                    f.append(i)
                    lis.append(t)
            except Exception as e:
                print(f"处理文件 {i} 时发生错误: {str(e)}")
    
    if not lis or len(lis) == 0:
        print("输入列表为空，无法进行预测。")
        return
    
    try:
        vectorizer_path = os.path.join(tmp['modph'], config['nmod'], 'vectorizer.pkl')
        classifier_path = os.path.join(tmp['modph'], config['nmod'], 'classifier.pkl')
        
        if not os.path.exists(vectorizer_path) or not os.path.exists(classifier_path):
            print("模型文件缺失，请检查模型路径和文件完整性")
            return
            
        vectorizer = joblib.load(vectorizer_path)
        clf = joblib.load(classifier_path)
        
        prediction = clf.predict(vectorizer.transform(lis))
        
        for i in range(len(f)):
            try:
                t = os.path.join(config['save'], prediction[i])
                
                # 创建目标目录（如果不存在）
                if not os.path.exists(t):
                    os.makedirs(t)
                    
                src_path = os.path.join(config['collect'], f[i])
                dst_path = os.path.join(t, f[i])
                
                # 文件操作重试机制（最多3次）
                max_retries = 3
                retry_count = 0
                
                if config['sepmd'] == 1:  # 移动文件
                    while retry_count <= max_retries:
                        try:
                            os.rename(src_path, dst_path)
                            break
                        except Exception as e:
                            retry_count += 1
                            if retry_count > max_retries:
                                print(f"移动文件失败：{src_path} -> {dst_path}，错误信息：{str(e)}")
                                raise
                            print(f"第 {retry_count} 次重试移动文件：{src_path}")
                            
                elif config['sepmd'] == 2:  # 复制文件
                    copy_file(src_path, dst_path)
                    
            except Exception as e:
                print(f"处理文件 {f[i]} 时发生错误: {str(e)}")
                continue
                
    except Exception as e:
        print(f"加载模型或预测过程中发生错误: {str(e)}")


# 复制文件
def copy_file(src, dst):
    with open(src, 'rb') as fsrc:
        with open(dst, 'wb') as fdst:
            while True:
                buf = fsrc.read(1024)
                if not buf:
                    break
                fdst.write(buf)


# 将文件夹路径转换为类别与文件路径的对应数据
def fold2file(n):
    lis = {}
    for i in os.listdir(n):
        i2 = os.path.join(n, i)
        if not os.path.isdir(i2):
            continue
        lis[i] = []
        for dirp, dirn, filen in os.walk(i2):
            for j in filen:
                k = os.path.join(dirp, j)
                if k.split('.')[-1] in ['pptx', 'docx', 'pdf']:
                    lis[i].append(k)
    return lis


# 使用可视化文件夹选择
def select_directory(tip):
    root = Tk()
    root.withdraw()
    return filedialog.askdirectory(title=tip)


# 读取配置文件
def config_rd():
    global config
    if not os.path.exists(os.path.join(os.path.dirname(__file__), 'data')):
        os.makedirs(os.path.join(os.path.dirname(__file__), 'data'))
    cfgph = os.path.join(os.path.dirname(__file__), 'data', 'config.json')
    if os.path.exists(cfgph):
        with open(cfgph, 'r') as file:
            config = json.load(file)


# 保存配置文件
def config_sv():
    global config
    with open(os.path.join(os.path.dirname(__file__), 'data', 'config.json'), 'w') as file:
        json.dump(config, file)


# 修改配置
def change_config(r):
    global config
    if r == 'collect':
        tp = select_directory('请选择文档“收集”文件夹，取消则默认为桌面')
        if tp == '':
            tp = os.path.join(os.path.expanduser('~'), 'Desktop')
        config['collect'] = tp.replace('/','\\')
    elif r == 'save':
        tp = select_directory('请选择文档“存放”文件夹，取消则默认为“收集”文件夹内的“已分类文档”文件夹')
        if tp == '':
            tp = os.path.join(config['collect'], '已分类文档')
            if not os.path.exists(tp):
                os.makedirs(tp)
        config['save'] = tp.replace('/','\\')
    elif r == 'nmod':
        nm = []
        for i in os.listdir(tmp['modph']):
            if os.path.isdir(os.path.join(tmp['modph'], i)):
                nm.append(i)
        while True:
            for i in range(len(nm)):
                print(f'{i + 1}.', nm[i])
            try:
                ip = int(input('请选择模型：'))
                t = os.path.join(tmp['modph'], nm[ip - 1])
                if os.path.exists(os.path.join(t, 'vectorizer.pkl')) and os.path.exists(
                        os.path.join(t, 'classifier.pkl')):
                    config['nmod'] = nm[ip - 1]
                    break
                print('没有此模型')
            except:
                print('error')
    elif r == 'sepmd':
        while True:
            try:
                print('\n请选择处理方式：1.移动 2.复制')
                ip = int(input('请输入：'))
                if ip not in [1, 2]:
                    print('error')
                    continue
                config['sepmd'] = ip
                break
            except:
                print('error')
    config_sv()


# 程序初始化
def initialize():
    config_rd()
    tmp['modph'] = os.path.join(os.path.dirname(__file__), 'data', 'model')
    if not os.path.exists(tmp['modph']):
        os.makedirs(tmp['modph'])
    if config['collect'] == '' or not os.path.exists(config['collect']):
        change_config('collect')
    if config['save'] == '' or not os.path.exists(config['save']):
        change_config('save')
    t = os.path.join(tmp['modph'], config['nmod'])
    if not (os.path.exists(os.path.join(t, 'vectorizer.pkl')) and os.path.exists(os.path.join(t, 'classifier.pkl'))):
        config['nmod'] = ''
        config_sv()


# 设置程序开机自启
def set_autostart():
    try:
        # 获取当前脚本的绝对路径
        current_path = os.path.abspath(sys.argv[0])

        # 如果是打包后的exe文件，获取实际路径
        if getattr(sys, 'frozen', False):
            current_path = sys.executable
        elif __file__:
            current_path = os.path.abspath(__file__)

        # 获取注册表启动项的键值
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, r"Software\Microsoft\Windows\CurrentVersion\Run", 0,
                             winreg.KEY_SET_VALUE)

        # 设置开机自启（将程序名称和路径写入注册表）
        winreg.SetValueEx(key, "DocumentClassifier", 0, winreg.REG_SZ, f'"{current_path}" -cl')
        winreg.CloseKey(key)
        return True
    except Exception as e:
        print(f"设置开机自启失败：{str(e)}")
        return False


# 取消程序开机自启
def unset_autostart():
    try:
        # 打开注册表启动项的键
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, r"Software\Microsoft\Windows\CurrentVersion\Run", 0,
                             winreg.KEY_SET_VALUE)

        # 删除自启项
        winreg.DeleteValue(key, "DocumentClassifier")
        winreg.CloseKey(key)
        return True
    except FileNotFoundError:
        # 如果键不存在，说明已经没有设置自启
        return True
    except Exception as e:
        print(f"取消开机自启失败：{str(e)}")
        return False


# 检查程序是否已设置为开机自启
def check_autostart():
    try:
        # 打开注册表启动项的键
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, r"Software\Microsoft\Windows\CurrentVersion\Run", 0,
                             winreg.KEY_READ)

        # 尝试读取值
        i = 0
        while True:
            name, value, type = winreg.EnumValue(key, i)
            if name == "DocumentClassifier":
                winreg.CloseKey(key)
                return True
            i += 1
    except OSError:
        # 当枚举完所有值或发生错误时退出
        pass
    finally:
        try:
            winreg.CloseKey(key)
        except:
            pass
    return False

#自动从github获取最新版本信息
def check_latest_version():
    global latest_version
    try:
        response = requests.get(
            'https://api.github.com/repos/ying-ck/Document-classifier/releases/latest',
            timeout=10
        )
        response.raise_for_status()  # 检查请求是否成功
        latest_version = response.json()['tag_name']

        skipped_version = None
        if os.path.exists('data'):
            skip_file = os.path.join('data', 'skip_update.txt')
            if os.path.exists(skip_file):
                with open(skip_file, 'r') as f:
                    skipped_version = f.read().strip()
        
        if version.parse(latest_version) > version.parse(ver) and latest_version != skipped_version:
            dialog = messagebox.askyesnocancel("新版本可用", "检测到新版本，是否前往GitHub下载？\n点击'是'将打开项目页面\n点击'否'将不再提示此版本更新\n点击'取消'跳过本次操作")
            
            if dialog is True:
                os.startfile("https://github.com/ying-ck/Document-classifier/releases")
            elif dialog is False:
                skip_file = os.path.join('data', 'skip_update.txt')
                with open(skip_file, 'w') as f:
                    f.write(latest_version)
                    
    except (requests.RequestException, KeyError) as e:
        print(f"获取版本信息失败: {e}")
        latest_version = "未知版本"


# 外部数据与必要设置
ver = 'v1.1.1'
jieba.set_dictionary(os.path.join(os.path.dirname(__file__), 'data', 'dict.txt'))
config = {'nmod': 'default', 'collect': '', 'save': '', 'sepmd': 1}
tmp = {}
initialize()
latest_version = None
version_thread = threading.Thread(target=check_latest_version)
version_thread.daemon = True
version_thread.start()

# 可视化界面相关代码
class DocumentClassifierGUI:
    def __init__(self, master):
        self.master = master
        master.title("Document-classifier " + ver)
        
        # 设置窗口大小和位置
        window_width = 600
        window_height = 400
        screen_width = master.winfo_screenwidth()
        screen_height = master.winfo_screenheight()
        x = (screen_width/2) - (window_width/2)
        y = (screen_height/2) - (window_height/2)
        master.geometry('%dx%d+%d+%d' % (window_width, window_height, x, y))
        
        # 创建主框架
        self.main_frame = Frame(master)
        self.main_frame.pack(pady=20, padx=20)
        
        # 创建标题标签
        self.title_label = Label(self.main_frame, text="Document-classifier", font=("Arial", 20, "bold"))
        self.title_label.pack(pady=10)
        
        # 创建作者信息标签
        self.author_label = Label(self.main_frame, text="作者：Yck", font=("Arial", 12))
        self.author_label.pack(pady=5)
        
        # GitHub链接
        self.github_label = Label(self.main_frame, text="GitHub项目地址：https://github.com/ying-ck/Document-classifier", 
                                 font=("Arial", 10), fg="blue", cursor="hand2")
        self.github_label.pack(pady=5)
        
        # 绑定点击事件到超链接
        self.github_label.bind("<Button-1>", lambda e: os.startfile("https://github.com/ying-ck/Document-classifier"))
        
        # 功能按钮框架
        self.function_frame = Frame(self.main_frame)
        self.function_frame.pack(pady=20)
        
        # 训练模型按钮
        self.train_button = Button(self.function_frame, text="训练模型", width=15, height=2, command=self.train_model)
        self.train_button.grid(row=0, column=0, padx=10, pady=10)
        
        # 进行分类按钮
        self.classify_button = Button(self.function_frame, text="进行分类", width=15, height=2, command=self.classify_documents)
        self.classify_button.grid(row=0, column=1, padx=10, pady=10)
        
        # 设置按钮
        self.settings_button = Button(self.function_frame, text="设置", width=15, height=2, command=self.show_settings)
        self.settings_button.grid(row=1, column=0, padx=10, pady=10)
        
        # 退出按钮
        self.exit_button = Button(self.function_frame, text="退出", width=15, height=2, command=master.quit)
        self.exit_button.grid(row=1, column=1, padx=10, pady=10)
        
        # 版权声明
        self.copyright_label = Label(self.main_frame, text="本程序完全免费", font=("Arial", 10))
        self.copyright_label.pack(pady=5)
    
    def train_model(self):
        folder = select_directory('请选择包含分类子文件夹的文件夹')
        if folder:
            model_sv(fold2file(folder))
            messagebox.showinfo("完成", "模型训练和保存已完成！")
    
    def classify_documents(self):
        model_path = os.path.join(tmp['modph'], config['nmod'])
        if config['nmod'] == '' or not (os.path.exists(os.path.join(model_path, 'vectorizer.pkl')) and 
                                        os.path.exists(os.path.join(model_path, 'classifier.pkl'))):
            change_config('nmod')
        
        progress_window = messagebox.showinfo("处理中", "正在分类文档，请稍候...")
        model_use()
        messagebox.showinfo("完成", "文档分类已完成！")
    
    def show_settings(self):
        SettingsWindow(self.master)
    

class SettingsWindow:
    def __init__(self, parent):
        self.parent = parent  # 保存父窗口引用
        self.window = Toplevel(parent)
        self.window.title("设置")
        
        # 设置窗口大小和位置
        window_width = 500
        window_height = 325
        screen_width = parent.winfo_screenwidth()
        screen_height = parent.winfo_screenheight()
        x = (screen_width/2) - (window_width/2)
        y = (screen_height/2) - (window_height/2)
        self.window.geometry('%dx%d+%d+%d' % (window_width, window_height, x, y))
        
        # 收集位置
        self.collect_frame = Frame(self.window)
        self.collect_frame.pack(pady=10, padx=20, fill="x")
        self.collect_label = Label(self.collect_frame, text="文档收集位置:")
        self.collect_label.pack(side="left")
        self.collect_var = StringVar(value=config['collect'] if config['collect'] else "未设置")
        self.collect_entry = Entry(self.collect_frame, textvariable=self.collect_var, width=40 ,state="readonly")
        self.collect_entry.pack(side="left", expand=True, fill="x")
        self.collect_button = Button(self.collect_frame, text="选择", command=self.change_collect)
        self.collect_button.pack(side="left", padx=(5, 0))
        
        # 存放位置
        self.save_frame = Frame(self.window)
        self.save_frame.pack(pady=10, padx=20, fill="x")
        self.save_label = Label(self.save_frame, text="文档存放位置:")
        self.save_label.pack(side="left")
        self.save_var = StringVar(value=config['save'] if config['save'] else "未设置")
        self.save_entry = Entry(self.save_frame, textvariable=self.save_var, width=40 ,state="readonly")
        self.save_entry.pack(side="left", expand=True, fill="x")
        self.save_button = Button(self.save_frame, text="选择", command=self.change_save)
        self.save_button.pack(side="left", padx=(5, 0))
        
        # 选择模型
        self.model_frame = Frame(self.window)
        self.model_frame.pack(pady=10, padx=20, fill="x")
        self.model_label = Label(self.model_frame, text="当前模型:")
        self.model_label.pack(side="left")
        
        # 创建模型选择下拉菜单
        model_path = os.path.join(tmp['modph'])
        self.available_models = []
        if os.path.exists(model_path):
            for item in os.listdir(model_path):
                item_path = os.path.join(model_path, item)
                if os.path.isdir(item_path) and \
                   os.path.exists(os.path.join(item_path, 'vectorizer.pkl')) and \
                   os.path.exists(os.path.join(item_path, 'classifier.pkl')):
                    self.available_models.append(item)
        
        self.model_var = StringVar(value=config['nmod'] if config['nmod'] and config['nmod'] in self.available_models else 
                                  ("未选择" if not self.available_models else self.available_models[0]))
        
        # 使用OptionMenu替换Entry用于模型选择
        self.model_menu = OptionMenu(self.model_frame, self.model_var, *self.available_models)
        self.model_menu.pack(side="left", expand=True, fill="x")
        
        # 开机自启选项
        self.autostart_frame = Frame(self.window)
        self.autostart_frame.pack(pady=10, padx=20, fill="x")
        self.autostart_label = Label(self.autostart_frame, text="开机自动分类(Beta):")
        self.autostart_label.pack(side="left")
        
        # 获取当前自启动状态
        self.is_autostart = check_autostart()
        self.autostart_state = StringVar(value="已启用" if self.is_autostart else "未启用")
        
        self.autostart_status = Label(self.autostart_frame, textvariable=self.autostart_state, width=10)
        self.autostart_status.pack(side="left", padx=(5, 0))
        
        self.autostart_toggle = Button(self.autostart_frame, text="更改", width=10, command=self.toggle_autostart)
        self.autostart_toggle.pack(side="left", padx=(5, 0))
        
        # 处理方式
        self.sepmd_frame = Frame(self.window)
        self.sepmd_frame.pack(pady=10, padx=20, fill="x")
        self.sepmd_label = Label(self.sepmd_frame, text="文档处理方式:")
        self.sepmd_label.pack(side="left")
        self.sepmd_var = StringVar(value="移动" if config['sepmd'] == 1 else "复制")
        self.sepmd_menu = Button(self.sepmd_frame, textvariable=self.sepmd_var, width=10, command=self.change_sepmd)
        self.sepmd_menu.pack(side="left")
        
        # 保存按钮
        self.save_button = Button(self.window, text="保存模型选择", width=15, command=self.save_settings)
        self.save_button.pack(pady=20)
        
        # 设置窗口为模态对话框，保持在主窗口之上
        self.window.transient(parent)
        self.window.grab_set()
        parent.wait_window(self.window)
    
    def change_collect(self):
        change_config('collect')
        self.collect_var.set(config['collect'])
    
    def change_save(self):
        change_config('save')
        self.save_var.set(config['save'])
    
    def change_model(self):
        """这个方法现在不会被调用，保留是为了防止错误"""
        pass
    
    def change_sepmd(self):
        current = config['sepmd']
        config['sepmd'] = 2 if current == 1 else 1
        self.sepmd_var.set("移动" if config['sepmd'] == 1 else "复制")
    
    def toggle_autostart(self):
        """切换开机自启状态"""
        if self.is_autostart:
            success = unset_autostart()
            if success:
                self.is_autostart = False
                self.autostart_state.set("未启用")
                messagebox.showinfo("提示", "已关闭开机自启！")
        else:
            success = set_autostart()
            if success:
                self.is_autostart = True
                self.autostart_state.set("已启用")
                messagebox.showinfo("提示", "已设置开机自启！")
    
    def save_settings(self):
        # 如果没有可用模型且用户未选择，则提示
        if not self.available_models:
            messagebox.showerror("错误", "没有可选的模型，请先训练一个模型。")
            return
            
        # 更新模型配置
        selected_model = self.model_var.get()
        if selected_model != config['nmod']:
            config['nmod'] = selected_model
        
        config_sv()
        messagebox.showinfo("提示", "设置已保存！")
        self.window.destroy()


# 主程序入口
if __name__ == '__main__':
    if len(sys.argv) > 1:
        if sys.argv[1] == '--help' or sys.argv[1] == '-h':
            print('使用方法：python main.py [--help | -h] [--version | -v] [--config | -c] [--classify | -cl]')
            print('--help | -h：显示帮助信息')
            print('--version | -v：显示版本信息')
            print('--config | -c：修改配置')
            print('--classify | -cl：开始分类文档')
        elif sys.argv[1] == '--version' or sys.argv[1] == '-v':
            print(ver)
        elif sys.argv[1] == '--config' or sys.argv[1] == '-c':
            if sys.argv[2:]:
                change_config(sys.argv[2])
            else:
                change_config(input('请选择要修改的配置：collect 收集文件夹 save 存放文件夹 nmod 模型 sepmd 处理方式\n'))
        elif  sys.argv[1] == '--classify' or sys.argv[1] == '-cl':
            model_use()
        else:
            print('无效的参数！')
    else:
        root = Tk()
        app = DocumentClassifierGUI(root)
        root.mainloop()

